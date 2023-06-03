from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches

def convert_pdf_to_ppt(pdf_path, ppt_path):
    # Convert PDF pages to images
    images = convert_from_path(pdf_path)

    # Create a new PowerPoint presentation
    prs = Presentation()

    # Add each image as a slide in the presentation
    for image in images:
        slide_layout = prs.slide_layouts[1]  # Use the second slide layout (Title and Content)
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.add_picture(image, Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))

    # Save the PowerPoint presentation
    prs.save(ppt_path)


# Example usage
convert_pdf_to_ppt( 'pdfconverter/pdfs/Writeup G12.pdf' , 'pdfconverter/file.pptx' )
